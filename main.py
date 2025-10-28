import os
import uuid
import logging
from io import BytesIO

from fastapi import FastAPI, File, UploadFile
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docx import Document
from docx.shared import Pt, RGBColor as DocxRGB, Inches
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from pptx.oxml.ns import qn


import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
import tempfile
import unicodedata, re


# -------------------------
# Logging
# -------------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
logger = logging.getLogger(__name__)

# -------------------------
# FastAPI app
# -------------------------
app = FastAPI(title="PPTX -> Word Converter (Full with Shapes)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],    # restrict in prod if needed
    allow_methods=["*"],
    allow_headers=["*"],
    allow_credentials=True,
)

# -------------------------
# Helpers
# -------------------------

def sanitize_text(s: str) -> str:
    """
    Remove control characters that break python-docx (e.g. NUL, other C0 controls).
    Keep common whitespace (space, tab, newline).
    """
    if s is None:
        return ""
    # Replace NULL bytes and other C0 controls except \n, \r, \t
    clean_chars = []
    removed = 0
    for ch in s:
        code = ord(ch)
        if code == 0:
            removed += 1
            continue
        # allow printable + common whitespace
        if code >= 32 or ch in ("\n", "\r", "\t"):
            clean_chars.append(ch)
        else:
            removed += 1
    clean = "".join(clean_chars)
    if removed:
        logger.debug(f"sanitize_text: removed {removed} control chars")
    return clean


def safe_filename(name: str) -> str:
    name = unicodedata.normalize("NFKD", name)
    name = re.sub(r"[^\w\-_. ]", "_", name)
    return name.strip() or "converted"

EMU_PER_INCH = 914400
def emu_to_inches(emu): return emu / EMU_PER_INCH

def render_shape_to_image(shape):
    """Draw PPT shape (rect/ellipse/triangle) as an image with text."""
    width_in = emu_to_inches(shape.width)
    height_in = emu_to_inches(shape.height)
    img_w, img_h = int(width_in * 300), int(height_in * 300)
    img = Image.new("RGB", (img_w, img_h), "white")
    draw = ImageDraw.Draw(img)

    # fill color
    fill_color = "#cccccc"
    try:
        if shape.fill and shape.fill.fore_color.type == 1:
            rgb = shape.fill.fore_color.rgb
            fill_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except Exception:
        pass

    # draw border + fill
    draw.rectangle([0, 0, img_w - 2, img_h - 2], fill=fill_color, outline="black", width=5)

    # add shape text
    try:
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            font = ImageFont.load_default()
            draw.text((20, img_h / 2.5), text, fill="black", font=font)
    except Exception as e:
        logger.debug(f"Shape text draw failed: {e}")

    tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    img.save(tmp_img.name)
    return tmp_img.name, width_in, height_in

def add_page_border(doc):
    """Add a double-line border around each page."""
    for section in doc.sections:
        section.top_margin = Inches(0.8)
        section.bottom_margin = Inches(0.8)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
        pgBorders = parse_xml(r'''
        <w:pgBorders %s>
            <w:top w:val="double" w:sz="6" w:space="24" w:color="000000"/>
            <w:left w:val="double" w:sz="6" w:space="24" w:color="000000"/>
            <w:bottom w:val="double" w:sz="6" w:space="24" w:color="000000"/>
            <w:right w:val="double" w:sz="6" w:space="24" w:color="000000"/>
        </w:pgBorders>
        ''' % nsdecls('w'))
        section._sectPr.append(pgBorders)

def pptx_color_to_rgb(color_obj):
    try:
        if color_obj and color_obj.type == 1:  # RGB
            rgb = color_obj.rgb
            return DocxRGB(rgb[0], rgb[1], rgb[2])
    except Exception as e:
        logger.debug(f"pptx_color_to_rgb: {e}")
    return None

def safe_get_text(shape):
    """Return text from a shape if available in a safe way."""
    try:
        if hasattr(shape, "text"):
            return shape.text or ""
    except Exception as e:
        logger.debug(f"safe_get_text: {e}")
    return ""

def save_stream_to_file(stream_bytes: bytes, ext: str = "png"):
    """Save bytes to a temporary file and return path."""
    path = f"/tmp/{uuid.uuid4()}.{ext}"
    with open(path, "wb") as f:
        f.write(stream_bytes)
    return path

def render_chart_from_chart_data(chart):
    """Render chart via matplotlib using chart_data."""
    try:
        chart_data = chart.chart_data
    except Exception as e:
        raise RuntimeError(f"No chart_data available: {e}")

    # categories (x-axis)
    try:
        categories = list(chart_data.categories)
        categories = [str(c) for c in categories]
    except Exception:
        categories = None

    # series
    series_list = []
    try:
        for s in chart_data.series:
            label = s.name if hasattr(s, "name") else None
            values = list(s.values)
            series_list.append((label, values))
    except Exception:
        raise RuntimeError("Failed to read series from chart_data")

    if not series_list:
        raise RuntimeError("No series data found in chart_data")

    # Figure
    fig, ax = plt.subplots(figsize=(6, 4))
    try:
        if len(series_list) == 1:
            label, values = series_list[0]
            if "pie" in str(chart.chart_type).lower():
                ax.pie(values, labels=categories if categories else None, autopct="%1.1f%%")
            elif "bar" in str(chart.chart_type).lower() or categories is not None:
                ax.bar(categories if categories else list(range(len(values))), values, label=label)
                if label:
                    ax.legend()
                ax.set_xticklabels(categories if categories else [str(i) for i in range(len(values))], rotation=45, ha="right")
            else:
                ax.plot(categories if categories else list(range(len(values))), values, marker="o", label=label)
                ax.legend()
        else:
            x = range(len(series_list[0][1]))
            width = 0.8 / max(1, len(series_list))
            for idx, (label, values) in enumerate(series_list):
                pos = [xi + (idx - len(series_list)/2) * width + width/2 for xi in x]
                ax.bar(pos, values, width=width, label=(label or f"Series {idx+1}"))
            ax.set_xticks(x)
            if categories:
                ax.set_xticklabels(categories, rotation=45, ha="right")
            ax.legend()
    except Exception as e:
        plt.close(fig)
        raise RuntimeError(f"Failed to render chart via matplotlib: {e}")

    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf.read()

def extract_image_from_shape(shape):
    """Try to extract image from shape."""
    try:
        if hasattr(shape, "image") and shape.image is not None:
            return shape.image.blob
    except:
        pass
    try:
        blip = shape.element.xpath(".//a:blip")
        if blip:
            rId = blip[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            rel = shape.part.related_parts.get(rId)
            if rel:
                return rel.blob
    except:
        pass
    raise RuntimeError("No image found in shape")

def draw_shape_as_image(shape):
    """Draw shape (rect, ellipse, arrow) with fill, border, and text."""
    scale = 2
    width = int(shape.width.pt * scale)
    height = int(shape.height.pt * scale)

    img = Image.new("RGBA", (width, height), (255,255,255,0))
    draw = ImageDraw.Draw(img)

    # Fill color
    fill_color = (255,255,255,0)
    try:
        if shape.fill.type == 1:
            fc = shape.fill.fore_color.rgb
            fill_color = (fc[0], fc[1], fc[2], 255)
    except: pass

    # Border color
    border_color = (0,0,0,255)
    border_width = 2
    try:
        if shape.line.color.type == 1:
            lc = shape.line.color.rgb
            border_color = (lc[0], lc[1], lc[2], 255)
        border_width = int(shape.line.width.pt*scale)
    except: pass

    stype = shape.shape_type
    if stype in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.RECTANGLE, MSO_SHAPE_TYPE.ROUNDED_RECTANGLE]:
        draw.rectangle([(0,0),(width,height)], fill=fill_color, outline=border_color, width=border_width)
    elif stype in [MSO_SHAPE_TYPE.ELLIPSE]:
        draw.ellipse([(0,0),(width,height)], fill=fill_color, outline=border_color, width=border_width)
    elif stype in [MSO_SHAPE_TYPE.ARROW, MSO_SHAPE_TYPE.CALLOUT]:
        draw.polygon([(0, height//2),(width-10,0),(width,height//2),(width-10,height)], fill=fill_color, outline=border_color)

    # Text inside shape
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            try:
                run = shape.text_frame.paragraphs[0].runs[0]
                font_name = run.font.name or "Arial"
                font_size = int((run.font.size.pt if run.font.size else 14)*scale)
                color = run.font.color.rgb
                text_color = (color[0],color[1],color[2],255) if color else (0,0,0,255)
            except:
                font_name, font_size, text_color = "Arial", 14, (0,0,0,255)
            try:
                font = ImageFont.truetype(f"{font_name}.ttf", font_size)
            except:
                font = ImageFont.load_default()
            bbox = draw.textbbox((0,0), text, font=font)
            text_x = (width - (bbox[2]-bbox[0]))//2
            text_y = (height - (bbox[3]-bbox[1]))//2
            draw.text((text_x,text_y), text, fill=text_color, font=font)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

# -------------------------
# Main endpoint
# -------------------------
@app.post("/convert/")
async def convert(file: UploadFile = File(...),slide_separator: int = 0):
    try:
        original_filename = file.filename or "uploaded.pptx"
        safe_name = safe_filename(os.path.splitext(original_filename)[0])

        logger.info(f"Received file: {original_filename.encode('utf-8', 'ignore').decode()}")
        
        content = await file.read()
        pptx_path = f"/tmp/{uuid.uuid4()}.pptx"
        with open(pptx_path,"wb") as f: f.write(content)
        logger.info(f"Saved PPTX to {pptx_path}")

        prs = Presentation(pptx_path)
        doc = Document()

        # Detect default font
        default_font_name = "Utsaah"
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    default_font_name = run.font.name
                                    break
                            if default_font_name!="Utsaah": break
                except: continue
            if default_font_name!="Utsaah": break
        logger.info(f"Default font detected: {default_font_name}")

        # Process slides
        for s_i, slide in enumerate(prs.slides, start=1):
            logger.info(f"Processing slide {s_i}/{len(prs.slides)}")

            for sh_i, shape in enumerate(slide.shapes, start=1):
                logger.info(f"Shape {sh_i}/{len(slide.shapes)} type={shape.shape_type}")

                # Text
                # text = safe_get_text(shape).strip()
                # if text and hasattr(shape, "text_frame") and shape.has_text_frame:
                #     try:
                #         for para in shape.text_frame.paragraphs:
                #             if not para.text.strip(): continue
                #             p = doc.add_paragraph()
                #             for run in para.runs:
                #                 r = p.add_run(run.text)
                #                 try:
                #                     r.font.name = default_font_name
                #                     r.font.size = Pt(14)
                #                     r.bold = run.font.bold
                #                     r.italic = run.font.italic
                #                     r.underline = run.font.underline
                #                     rgb = pptx_color_to_rgb(run.font.color)
                #                     if rgb: r.font.color.rgb = rgb
                #                 except: pass
                #     except:
                #         doc.add_paragraph(text)

                # 
                # Text (with robust bullets + sanitization
                try:
                    if hasattr(shape, "text_frame") and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            raw_para_text = para.text or ""
                            if not raw_para_text.strip():
                                continue
                
                            # --- Detect bullet or numbering type safely ---
                            list_style = None
                            try:
                                pPr = getattr(para._p, "pPr", None)
                                if pPr is not None:
                                    # Regular bullet (•, ○, –, →)
                                    buChar = pPr.find(qn("a:buChar"))
                                    if buChar is not None:
                                        list_style = "List Bullet"
                
                                    # Numbered list (1., i., etc.)
                                    elif pPr.find(qn("a:buAutoNum")) is not None:
                                        list_style = "List Number"
                
                                    # No bullet explicitly set
                                    elif pPr.find(qn("a:buNone")) is not None:
                                        list_style = None
                            except Exception as e:
                                logger.debug(f"Bullet detection failed: {e}")
                
                            # --- Choose Word paragraph style ---
                            if list_style == "List Bullet":
                                p = doc.add_paragraph(style="List Bullet")
                            elif list_style == "List Number":
                                p = doc.add_paragraph(style="List Number")
                            else:
                                p = doc.add_paragraph()
                
                            # --- Apply indentation level (PowerPoint nesting) ---
                            try:
                                level = getattr(para, "level", 0)
                                if isinstance(level, int) and level > 0:
                                    # Each level adds 0.5 cm indentation
                                    p.paragraph_format.left_indent = Inches(0.2 * level)
                            except Exception as e:
                                logger.debug(f"Indentation failed: {e}")
                
                            # --- Add text runs (with styling + sanitization) ---
                            for run in para.runs:
                                run_text = run.text or ""
                                run_text = sanitize_text(run_text)
                                if not run_text.strip():
                                    continue
                
                                try:
                                    r = p.add_run(run_text)
                                    # Preserve styling
                                    try:
                                        r.font.name = default_font_name
                                        r.font.size = Pt(14)
                                        r.bold = bool(run.font.bold)
                                        r.italic = bool(run.font.italic)
                                        r.underline = bool(run.font.underline)
                                        rgb = pptx_color_to_rgb(run.font.color)
                                        if rgb:
                                            r.font.color.rgb = rgb
                                    except Exception as style_err:
                                        logger.debug(f"Run styling skipped: {style_err}")
                
                                except Exception as e:
                                    # fallback if font fails
                                    try:
                                        p.add_run(sanitize_text(run_text))
                                    except Exception as e2:
                                        logger.error(f"Failed to add sanitized run: {e2}")
                
                except Exception as e:
                    logger.warning(f"Text extraction failed (fallback): {e}")
                    try:
                        fallback_text = sanitize_text(safe_get_text(shape))
                        if fallback_text:
                            doc.add_paragraph(fallback_text)
                    except Exception as e2:
                        logger.error(f"Failed fallback text insertion: {e2}")

                


                # Table
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        word_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                        word_table.style = "Table Grid"
                        for r_idx, row in enumerate(table.rows):
                            for c_idx, cell in enumerate(row.cells):
                                txt = cell.text.strip()
                                p = word_table.cell(r_idx, c_idx).paragraphs[0]
                                run = p.runs[0] if p.runs else p.add_run(txt)
                                run.font.size = Pt(14)
                                run.font.name = default_font_name

                except: pass

                # # Image
                # try:
                #     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape,"image") or "blip" in str(shape.element.xml):
                #         try:
                #             img_bytes = extract_image_from_shape(shape)
                #             doc.add_picture(BytesIO(img_bytes), width=Inches(4))
                #         except: pass
                # except: pass
                # Image
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image") or "blip" in str(shape.element.xml):
                        try:
                            img_bytes = extract_image_from_shape(shape)
                            width_in = emu_to_inches(shape.width)
                            height_in = emu_to_inches(shape.height)
                            doc.add_picture(BytesIO(img_bytes), width=Inches(width_in), height=Inches(height_in))
                            logger.info(f"Inserted image with original size: {width_in:.2f} x {height_in:.2f} inches")
                        except Exception as e:
                            logger.warning(f"Image extraction failed: {e}")
                except Exception as e:
                    logger.debug(f"Image block error: {e}")

                # Chart
                try:
                    if hasattr(shape,"chart"):
                        chart_inserted=False
                        try:
                            chart_part = getattr(shape.chart,"chart_part",None)
                            if chart_part:
                                img_stream = BytesIO(chart_part.chart_space.blob)
                                doc.add_picture(img_stream, width=Inches(5))
                                chart_inserted=True
                        except: pass
                        if not chart_inserted:
                            try:
                                img_bytes = render_chart_from_chart_data(shape.chart)
                                doc.add_picture(BytesIO(img_bytes), width=Inches(5))
                                chart_inserted=True
                            except:
                                doc.add_paragraph("[Chart could not be rendered]")
                except: pass

                # Shapes
                try:
                    if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE,
                                            MSO_SHAPE_TYPE.ARROW, MSO_SHAPE_TYPE.CALLOUT,
                                            MSO_SHAPE_TYPE.ROUNDED_RECTANGLE, MSO_SHAPE_TYPE.RECTANGLE]:
                        img_path, w_in, h_in = render_shape_to_image(shape)
                        doc.add_picture(img_path, width=Inches(w_in), height=Inches(h_in))
                        logger.info(f"Rendered and inserted shape: {shape.shape_type}")
                except Exception as e:
                    logger.warning(f"Shape render failed: {e}")

            # Slide separation: 2 blank lines
            # -------------------------
            # Handle slide separation logic
            # -------------------------
            if slide_separator == -1:
                # New page for each slide
                doc.add_page_break()
            elif slide_separator == 0:
                # No extra spacing between slides
                pass
            else:
                # Add N blank lines between slides
                for _ in range(slide_separator):
                    doc.add_paragraph("")


        # Save Word
        out_path = f"/tmp/{uuid.uuid4()}.docx"
        add_page_border(doc)
        doc.save(out_path)
        logger.info(f"Saved Word to {out_path}")

        def iterfile():
            with open(out_path,"rb") as f: yield from f

        return StreamingResponse(
            iterfile(),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": f'attachment; filename="{safe_name}.docx"; filename*=UTF-8\'\'{safe_name}.docx'
            },
        )
    except Exception as e:
        logger.exception("Conversion failed")
        return JSONResponse({"error": str(e)}, status_code=500)

# Health check
@app.get("/health")
def health(): return {"status":"ok"}

# Path to your frontend file
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
INDEX_PATH = os.path.join(BASE_DIR, "index.html")

@app.get("/")
def serve_index():
    if os.path.exists(INDEX_PATH):
        return FileResponse(INDEX_PATH, media_type="text/html")
    else:
        return {"error": "index.html not found"}
